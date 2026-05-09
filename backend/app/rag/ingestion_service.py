"""Document ingestion service for RAG platform"""

import asyncio
import logging
import io
import re
from typing import Optional
from uuid import UUID
from datetime import datetime, timezone

import pdfplumber
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation
from sqlalchemy.ext.asyncio import AsyncSession
from sqlalchemy import select, update

from app.core.config import settings
from app.database.core import AsyncSessionLocal
from app.entities.document import Document, IngestionStatus
from app.notifications.redis_notifier import redis_notifier
from app.rag.chunker import Chunker
from app.rag.embedding_config import build_embeddings_client
from app.rag.vector_store import VectorStore

logger = logging.getLogger(__name__)

_POSTGRES_UNSUPPORTED_CHARS = re.compile(r"\x00")


def sanitize_text_for_postgres(text: str) -> str:
    """Remove characters PostgreSQL text columns cannot store."""

    return _POSTGRES_UNSUPPORTED_CHARS.sub("", text)


class IngestionService:
    """Handles document ingestion pipeline: extract -> chunk -> embed -> store"""

    def __init__(self):
        self.chunker = Chunker()
        self.vector_store = VectorStore()
        self._embeddings = self._build_embeddings()

    @staticmethod
    def _build_embeddings():
        return build_embeddings_client()

    async def ingest_document(
        self,
        document_id: UUID,
        user_id: UUID,
        file_content: bytes,
        filename: str,
    ) -> dict:
        """
        Full ingestion pipeline for a document.

        Args:
            document_id: Document UUID
            user_id: User UUID
            file_content: Document file bytes (PDF, DOCX, XLSX, XLS, PPTX, TXT, MD, CSV)
            filename: Original filename

        Returns:
            Dict with status and metrics
        """
        async with AsyncSessionLocal() as db:
            try:
                logger.info(f"Starting ingestion for document {document_id}")
                # Update document status to processing
                await self._update_document_status(
                    db, document_id, IngestionStatus.PROCESSING, 0
                )
                await db.commit()
                await redis_notifier.publish(
                    user_id,
                    {
                        "type": "document_status",
                        "data": {
                            "document_id": str(document_id),
                            "status": "processing",
                            "progress": 0,
                        },
                    },
                )

                # Step 1: Extract text from document
                text = await self._extract_text(file_content, filename)
                text = sanitize_text_for_postgres(text)
                if not text:
                    raise ValueError("No text extracted from document")

                # Step 2: Chunk text
                logger.info("Chunking text")
                chunks_data = self.chunker.create_chunks(
                    text, metadata={"filename": filename}
                )

                total_chunks = len(chunks_data)

                # Validate chunks were generated
                if total_chunks == 0:
                    logger.error(f"Document {document_id} produced 0 chunks")
                    await self._update_document_status(
                        db,
                        document_id,
                        IngestionStatus.FAILED,
                        0,
                        error_message="El documento no tiene suficiente contenido para ser procesado (0 chunks generados)",
                    )
                    await db.commit()
                    raise ValueError(
                        "Document has insufficient content (0 chunks generated)"
                    )

                await self._update_document_status(
                    db, document_id, IngestionStatus.PROCESSING, 10
                )
                await db.commit()

                # Step 3: Generate embeddings
                logger.info(f"Generating embeddings for {total_chunks} chunks")
                embeddings = await self._generate_embeddings(
                    [chunk["content"] for chunk in chunks_data]
                )

                # Step 4: Store chunks with embeddings
                logger.info("Storing chunks in database")
                chunk_tuples = [
                    (
                        document_id,
                        user_id,
                        chunk["content"],
                        embedding,
                        chunk["chunk_index"],
                        chunk["token_count"],
                        chunk.get("metadata", {}),
                    )
                    for chunk, embedding in zip(chunks_data, embeddings)
                ]

                await self.vector_store.add_chunks(db, chunk_tuples)

                # Update document as completed
                await self._update_document(
                    db,
                    document_id,
                    {
                        "status": IngestionStatus.COMPLETED,
                        "total_chunks": total_chunks,
                        "processed_chunks": total_chunks,
                        "processed_at": datetime.now(timezone.utc),
                    },
                )
                await db.commit()

                await redis_notifier.publish(
                    user_id,
                    {
                        "type": "document_status",
                        "data": {
                            "document_id": str(document_id),
                            "status": "completed",
                            "progress": 100,
                        },
                    },
                )
                logger.info(f"Document ingestion completed: {document_id}")
                return {
                    "status": "completed",
                    "total_chunks": total_chunks,
                    "total_tokens": sum(chunk["token_count"] for chunk in chunks_data),
                }

            except Exception as e:
                error_message = self._safe_ingestion_error_message(e)
                logger.exception("Error ingesting document %s", document_id)

                # A DB error during chunk insertion leaves the transaction aborted.
                # Roll it back before trying to persist the FAILED status; otherwise
                # the document remains at the last committed PROCESSING state.
                await db.rollback()
                await self._update_document_status(
                    db,
                    document_id,
                    IngestionStatus.FAILED,
                    0,
                    error_message=error_message,
                )
                await db.commit()
                await redis_notifier.publish(
                    user_id,
                    {
                        "type": "document_status",
                        "data": {
                            "document_id": str(document_id),
                            "status": "failed",
                            "progress": 0,
                            "error_message": error_message,
                        },
                    },
                )
                return {"status": "failed", "error_message": error_message}

    @staticmethod
    def _safe_ingestion_error_message(error: Exception) -> str:
        """Return a useful ingestion error without leaking SQL parameters/content."""

        raw_message = str(error) or error.__class__.__name__
        lowered = raw_message.lower()
        if "invalid byte sequence for encoding" in lowered and "0x00" in lowered:
            return (
                "Document text contains unsupported NUL characters. "
                "The extracted text must be sanitized before storing chunks."
            )

        if error.__class__.__name__ in {"DBAPIError", "StatementError"}:
            return "Database error while storing document chunks. Check backend logs for details."

        return raw_message[:1000]

    async def _extract_text(self, file_content: bytes, filename: str) -> str:
        """Extract text from file based on extension"""
        ext = filename.lower().split(".")[-1]
        logger.info(f"Extracting text from {ext} file: {filename}")

        try:
            if ext == "pdf":
                return await self._extract_pdf_text(file_content)
            elif ext == "docx":
                return await self._extract_docx_text(file_content)
            elif ext in ("xlsx", "xls"):
                return await self._extract_excel_text(file_content)
            elif ext == "pptx":
                return await self._extract_pptx_text(file_content)
            elif ext in ("txt", "md", "csv"):
                return await self._extract_text_plain(file_content)
            else:
                raise ValueError(f"Unsupported file format: {ext}")
        except Exception as e:
            logger.error(f"Error extracting text from {filename}: {e}")
            raise

    async def _extract_pdf_text(self, file_content: bytes) -> str:
        """Extract text from PDF using pdfplumber — offloaded to thread pool to avoid blocking the event loop"""

        def _sync() -> str:
            text = ""
            with pdfplumber.open(io.BytesIO(file_content)) as pdf:
                for page in pdf.pages:
                    page_text = page.extract_text()
                    if page_text:
                        text += page_text + "\n\n"
            return text.strip()

        try:
            return await asyncio.to_thread(_sync)
        except Exception as e:
            logger.error(f"Error extracting PDF text: {e}")
            raise

    async def _extract_docx_text(self, file_content: bytes) -> str:
        """Extract text from Word document — offloaded to thread pool"""

        def _sync() -> str:
            doc = Document(io.BytesIO(file_content))
            text = ""
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    text += paragraph.text + "\n"
            return text.strip()

        try:
            return await asyncio.to_thread(_sync)
        except Exception as e:
            logger.error(f"Error extracting DOCX text: {e}")
            raise

    async def _extract_excel_text(self, file_content: bytes) -> str:
        """Extract text from Excel spreadsheet — offloaded to thread pool"""

        def _sync() -> str:
            workbook = load_workbook(
                io.BytesIO(file_content), read_only=True, data_only=True
            )
            text = ""
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                text += f"\n--- Sheet: {sheet_name} ---\n"
                for row in sheet.iter_rows(values_only=True):
                    row_text = " ".join(
                        str(cell) if cell is not None else "" for cell in row
                    )
                    if row_text.strip():
                        text += row_text + "\n"
            return text.strip()

        try:
            return await asyncio.to_thread(_sync)
        except Exception as e:
            logger.error(f"Error extracting Excel text: {e}")
            raise

    async def _extract_pptx_text(self, file_content: bytes) -> str:
        """Extract text from PowerPoint presentation — offloaded to thread pool"""

        def _sync() -> str:
            prs = Presentation(io.BytesIO(file_content))
            text = ""
            for i, slide in enumerate(prs.slides):
                text += f"\n--- Slide {i + 1} ---\n"
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text += shape.text + "\n"
            return text.strip()

        try:
            return await asyncio.to_thread(_sync)
        except Exception as e:
            logger.error(f"Error extracting PPTX text: {e}")
            raise

    async def _extract_text_plain(self, file_content: bytes) -> str:
        """Extract text from plain text or markdown file"""
        try:
            text = file_content.decode("utf-8")
            return text.strip()
        except UnicodeDecodeError:
            try:
                text = file_content.decode("latin-1")
                return text.strip()
            except Exception as e:
                logger.error(f"Error decoding text file: {e}")
                raise ValueError(
                    "Could not decode text file. Ensure it's UTF-8 or Latin-1 encoded."
                )
        except Exception as e:
            logger.error(f"Error extracting plain text: {e}")
            raise

    async def _generate_embeddings(
        self, texts: list[str], batch_size: int | None = None
    ) -> list[list[float]]:
        """
        Generate embeddings for document chunks.

        Args:
            texts: List of text strings
            batch_size: Number of texts per request

        Returns:
            List of embedding vectors
        """
        batch_size = batch_size or settings.EMBEDDING_BATCH_SIZE

        return await self._generate_openai_embeddings(texts, batch_size)

    @staticmethod
    def _is_rate_limit_error(error: Exception) -> bool:
        status_code = getattr(error, "status_code", None)
        if status_code == 429:
            return True

        error_str = str(error).lower()
        return (
            "429" in error_str
            or "rate limit" in error_str
            or "resource_exhausted" in error_str
        )

    @staticmethod
    def _external_embedding_error_message(provider: str, error: Exception) -> str:
        if IngestionService._is_rate_limit_error(error):
            provider_name = {
                "voyage": "Voyage AI",
                "zhipuai": "ZhipuAI",
                "openai": "OpenAI",
            }.get(provider, provider)
            return (
                f"{provider_name} rate limit exceeded while generating embeddings. "
                "Try again later, reduce EMBEDDING_BATCH_SIZE, or use a paid/higher-quota key."
            )
        return str(error) or error.__class__.__name__

    async def _retry_embedding_batch(self, provider: str, batch_num: int, operation):
        max_retries = settings.EMBEDDING_MAX_RETRIES
        for attempt in range(max_retries):
            try:
                return await operation()
            except Exception as e:
                if not self._is_rate_limit_error(e) or attempt == max_retries - 1:
                    message = self._external_embedding_error_message(provider, e)
                    raise RuntimeError(message) from e

                wait_time = min(
                    settings.EMBEDDING_RETRY_BASE_SECONDS * (2**attempt),
                    settings.EMBEDDING_RETRY_MAX_SECONDS,
                )
                logger.warning(
                    "%s embeddings rate limited on batch %s (attempt %s/%s). Waiting %.1fs...",
                    provider,
                    batch_num,
                    attempt + 1,
                    max_retries,
                    wait_time,
                )
                await asyncio.sleep(wait_time)

    async def _generate_openai_embeddings(
        self, texts: list[str], batch_size: int = 100
    ) -> list[list[float]]:
        """Generate embeddings using OpenAI-compatible API (Voyage/ZhipuAI/OpenAI)."""
        all_embeddings = []
        provider = settings.embedding_provider
        total_batches = (len(texts) + batch_size - 1) // batch_size

        for i in range(0, len(texts), batch_size):
            batch = texts[i : i + batch_size]
            batch_num = i // batch_size + 1
            logger.info(
                f"Generating {provider} embeddings for batch {batch_num}/{total_batches} "
                f"({len(batch)} texts)"
            )

            extra_body = {
                "input_type": "document",
                "output_dimension": settings.EMBEDDING_DIMENSION,
            }

            response = await self._retry_embedding_batch(
                provider,
                batch_num,
                lambda: self._embeddings.embeddings.create(
                    model=settings.embedding_model,
                    input=batch,
                    extra_body=extra_body,
                ),
            )

            batch_embeddings = [item.embedding for item in response.data]
            all_embeddings.extend(batch_embeddings)
            logger.info(
                f"Batch {batch_num}/{total_batches} complete: {len(batch_embeddings)} embeddings"
            )
        return all_embeddings

    async def _update_document_status(
        self,
        db: AsyncSession,
        document_id: UUID,
        status: IngestionStatus,
        progress: int,
        error_message: Optional[str] = None,
    ) -> None:
        """Update document status and progress"""
        await db.execute(
            update(Document)
            .where(Document.id == document_id)
            .values(
                status=status, processed_chunks=progress, error_message=error_message
            )
        )

    async def _update_document(
        self, db: AsyncSession, document_id: UUID, updates: dict
    ) -> None:
        """Update document with multiple fields"""
        await db.execute(
            update(Document).where(Document.id == document_id).values(**updates)
        )

    async def get_document(
        self, db: AsyncSession, document_id: UUID, user_id: UUID
    ) -> Optional[Document]:
        """Get document by ID with user isolation"""
        result = await db.execute(
            select(Document)
            .where(Document.id == document_id)
            .where(Document.user_id == user_id)
        )
        return result.scalar_one_or_none()
